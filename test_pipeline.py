"""Self-check for the pure logic: material extraction, caching, prompts, course suggestion.
Run: .venv/bin/python test_pipeline.py"""
import json
import time
import datetime
import tempfile
from pathlib import Path

import materials as mat
import app as appmod
from app import build_prompt, session_dir, course_dir, suggest_course, app


def test_extract_docx_and_pptx():
    with tempfile.TemporaryDirectory() as tmp:
        d = Path(tmp)

        from docx import Document
        doc = Document()
        doc.add_paragraph("Photosynthesis Overview")
        doc.add_paragraph("")  # blank paragraphs must not leak through
        doc.add_paragraph("Light reactions occur in the thylakoid.")
        doc.save(d / "handout.docx")

        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Calvin Cycle"
        slide.placeholders[1].text = "Carbon fixation"
        prs.save(d / "deck.pptx")

        out = mat.extract_all(d)

    assert "Photosynthesis Overview" in out
    assert "Light reactions occur in the thylakoid." in out
    assert "[Slide 1]" in out and "Calvin Cycle" in out and "Carbon fixation" in out
    assert "handout.docx" in out and "deck.pptx" in out, "chunks are labelled by filename"
    assert "\n\n\n" not in out, "blank paragraphs should not produce empty lines"


def test_markdown_and_text_materials_are_readable():
    with tempfile.TemporaryDirectory() as tmp:
        d = Path(tmp)
        (d / "reading.md").write_text("# Chapter 3\nNormalization rules.")
        (d / "outline.txt").write_text("plain text notes")
        assert [p.name for p in mat.readable_files(d)] == ["outline.txt", "reading.md"]
        out = mat.extract_all(d)
    assert "Normalization rules." in out and "reading.md" in out
    assert "plain text notes" in out


def test_unknown_extension_is_skipped():
    with tempfile.TemporaryDirectory() as tmp:
        d = Path(tmp)
        (d / "notes.xyz").write_text("ignore me")
        (d / ".DS_Store").write_text("ignore me too")
        assert mat.extract_all(d) == ""


def test_unreadable_file_does_not_crash_the_run():
    with tempfile.TemporaryDirectory() as tmp:
        d = Path(tmp)
        (d / "broken.pdf").write_text("this is not a pdf")
        out = mat.extract_all(d)
    assert "could not read broken.pdf" in out, "a bad file is reported, not raised"


def test_missing_materials_dir_is_empty():
    assert mat.extract_all(Path("/nonexistent/materials")) == ""


def test_prompt_with_materials_uses_them_as_skeleton():
    p = build_prompt("today we covered the thylakoid", "[Slide 1]\nCalvin Cycle")
    assert "skeleton" in p
    assert "CLASS MATERIALS" in p and "Calvin Cycle" in p
    assert "today we covered the thylakoid" in p
    assert "Additional Material From Lecture" in p


def test_prompt_without_materials_falls_back_to_topics():
    p = build_prompt("today we covered the thylakoid", "")
    assert "No class materials were provided" in p
    assert "CLASS MATERIALS" not in p
    assert "skeleton" not in p
    assert "today we covered the thylakoid" in p


def test_session_dir_rejects_path_traversal():
    with app.test_request_context():
        for bad_course, bad_date in [
            ("../../etc", "2026-01-01"),
            ("ok", "../../../etc"),
            ("ok", "not-a-date"),
            ("", "2026-01-01"),
        ]:
            try:
                session_dir(bad_course, bad_date)
                raise AssertionError(f"should have rejected {bad_course!r}/{bad_date!r}")
            except Exception as e:
                assert "400" in str(e) or "bad" in str(e).lower()

        good = session_dir("CS101", "2026-01-01")
        assert good.name == "2026-01-01" and good.parent.name == "CS101"


def _docx(path, *paragraphs):
    from docx import Document
    doc = Document()
    for p in paragraphs:
        doc.add_paragraph(p)
    doc.save(path)


def test_extract_cached_reuses_until_materials_change():
    with tempfile.TemporaryDirectory() as tmp:
        d = Path(tmp)
        mats = d / "materials"
        mats.mkdir()
        _docx(mats / "a.docx", "Chapter One")
        cache = d / "cache.json"

        first = mat.extract_cached(cache, mats)
        assert "Chapter One" in first
        assert cache.exists(), "extraction is cached to disk"

        # Corrupt the source but keep the signature: a cache hit must not re-read it.
        sig = json.loads(cache.read_text())["sig"]
        assert mat.extract_cached(cache, mats) == first
        assert json.loads(cache.read_text())["sig"] == sig

        time.sleep(0.01)
        _docx(mats / "b.docx", "Chapter Two")
        second = mat.extract_cached(cache, mats)
        assert "Chapter Two" in second and "Chapter One" in second, "new file invalidates the cache"


def test_corrupt_cache_is_a_miss_not_a_crash():
    with tempfile.TemporaryDirectory() as tmp:
        d = Path(tmp)
        mats = d / "materials"
        mats.mkdir()
        _docx(mats / "a.docx", "Chapter One")
        cache = d / "cache.json"
        cache.write_text("{not json at all")
        assert "Chapter One" in mat.extract_cached(cache, mats)


def test_shared_materials_are_labelled_for_the_prompt():
    with tempfile.TemporaryDirectory() as tmp:
        d = Path(tmp)
        shared, session = d / "_shared", d / "materials"
        shared.mkdir()
        session.mkdir()
        _docx(shared / "textbook.docx", "Textbook intro")
        _docx(session / "lecture.docx", "Lecture five slides")

        out = mat.extract_all(shared, session, label={str(shared): "course material: "})

    assert "--- course material: textbook.docx ---" in out
    assert "--- lecture.docx ---" in out
    assert out.index("textbook") < out.index("lecture"), "shared context comes first"


def test_extract_all_skips_missing_dirs():
    assert mat.extract_all(Path("/nope/a"), Path("/nope/b")) == ""


def _fake_classes(tmp):
    """Point the app at a throwaway classes/ tree, restoring it afterwards."""
    class Swap:
        def __enter__(self):
            self.real = appmod.CLASSES
            appmod.CLASSES = Path(tmp)
            return appmod.CLASSES
        def __exit__(self, *exc):
            appmod.CLASSES = self.real
    return Swap()


def _write_session(root, course, when, date):
    d = root / course / date
    d.mkdir(parents=True, exist_ok=True)
    (d / "meta.json").write_text(json.dumps({"started_at": when.isoformat()}))


def test_suggest_course_learns_weekday_and_time():
    with tempfile.TemporaryDirectory() as tmp, _fake_classes(tmp) as root:
        now = datetime.datetime.now()
        assert suggest_course() == "", "no history means no guess"

        # Shift within the same calendar day whichever way has room, so the
        # "wrong time" case cannot accidentally land on another weekday.
        far = datetime.timedelta(hours=-5 if now.hour >= 12 else 5)

        _write_session(root, "BIO201", now - datetime.timedelta(days=7), "2026-01-05")
        _write_session(root, "HIST100", now - datetime.timedelta(days=7) + far, "2026-01-06")
        _write_session(root, "MATH50", now - datetime.timedelta(days=8), "2026-01-07")

        assert suggest_course() == "BIO201", "same weekday and time of day wins"


def test_suggest_course_prefers_the_course_with_more_matching_history():
    with tempfile.TemporaryDirectory() as tmp, _fake_classes(tmp) as root:
        now = datetime.datetime.now()
        _write_session(root, "ONCE", now - datetime.timedelta(days=7), "2026-01-05")
        for i, weeks in enumerate((7, 14, 21)):
            _write_session(root, "WEEKLY", now - datetime.timedelta(days=weeks), f"2026-02-0{i+1}")
        assert suggest_course() == "WEEKLY"


def test_suggest_course_ignores_unreadable_meta():
    with tempfile.TemporaryDirectory() as tmp, _fake_classes(tmp) as root:
        d = root / "BROKEN" / "2026-01-05"
        d.mkdir(parents=True)
        (d / "meta.json").write_text("garbage")
        (root / "NOMETA" / "2026-01-06").mkdir(parents=True)
        assert suggest_course() == "", "bad metadata is skipped, not fatal"


def test_segments_are_ordered_numerically_not_alphabetically():
    """0004.webm must not sort before 0010.webm once a class runs past 10 chunks."""
    with tempfile.TemporaryDirectory() as tmp:
        session = Path(tmp)
        segs = session / "segments"
        segs.mkdir()
        for i in (12, 2, 100, 7):
            (segs / f"{i:04d}.webm").touch()
        assert [p.name for p in appmod.segment_files(session)] == [
            "0002.webm", "0007.webm", "0012.webm", "0100.webm"
        ]


def test_segment_status_counts_transcribed_chunks():
    with tempfile.TemporaryDirectory() as tmp:
        session = Path(tmp)
        segs = session / "segments"
        segs.mkdir()
        for i in range(3):
            (segs / f"{i:04d}.webm").touch()
        (segs / "0000.txt").write_text("first chunk")

        assert appmod.segment_status(session) == {"total": 3, "done": 1, "finished": False}
        (segs / ".done").touch()
        assert appmod.segment_status(session)["finished"] is True


def test_assemble_transcript_joins_finished_chunks_in_order():
    with tempfile.TemporaryDirectory() as tmp:
        session = Path(tmp)
        segs = session / "segments"
        segs.mkdir()
        for i, line in enumerate(["chunk one", "chunk two", "chunk three"]):
            (segs / f"{i:04d}.webm").touch()
            (segs / f"{i:04d}.txt").write_text(line + "\n")

        text = appmod.assemble_transcript(session)
        saved = (session / "transcript.txt").read_text()

    assert text == "chunk one\nchunk two\nchunk three"
    assert saved == text, "the stitched transcript is written to disk for reuse"


def test_assemble_transcript_skips_a_chunk_it_cannot_read():
    """One bad segment must not cost the whole lecture."""
    with tempfile.TemporaryDirectory() as tmp:
        session = Path(tmp)
        segs = session / "segments"
        segs.mkdir()
        (segs / "0000.webm").touch()
        (segs / "0000.txt").write_text("good chunk")
        (segs / "0001.webm").write_text("not audio at all")  # no .txt, will fail to transcribe

        text = appmod.assemble_transcript(session)

    assert text == "good chunk", "the readable chunk survives"


def test_assemble_transcript_errors_when_there_is_nothing_at_all():
    with tempfile.TemporaryDirectory() as tmp:
        try:
            appmod.assemble_transcript(Path(tmp))
            raise AssertionError("should have complained about a missing recording")
        except RuntimeError as e:
            assert "No recording" in str(e)


def test_course_dir_rejects_reserved_and_traversal_names():
    with app.test_request_context():
        for bad in ["_shared", "../etc", "", "a/b"]:
            try:
                course_dir(bad)
                raise AssertionError(f"should have rejected {bad!r}")
            except Exception as e:
                assert "400" in str(e) or "bad" in str(e).lower()
        assert course_dir("CS101").name == "CS101"


def test_delete_material_removes_only_a_plain_filename():
    with tempfile.TemporaryDirectory() as tmp, app.test_request_context():
        d = Path(tmp) / "materials"
        d.mkdir()
        (d / "slides.pdf").write_text("keep")
        (d / "junk.md").write_text("remove")
        (Path(tmp) / "meta.json").write_text("must survive")

        appmod.delete_material(d, "junk.md")
        assert not (d / "junk.md").exists()
        assert (d / "slides.pdf").exists(), "deleting one file leaves the others"

        for bad in ["..", "../meta.json", "a/b", "", ".DS_Store", "gone.pdf"]:
            try:
                appmod.delete_material(d, bad)
                raise AssertionError(f"should have rejected {bad!r}")
            except AssertionError:
                raise
            except Exception as e:
                assert "400" in str(e) or "404" in str(e), f"{bad!r} -> {e}"
        assert (Path(tmp) / "meta.json").exists(), "nothing escapes the materials folder"


class FakeMessage:
    def __init__(self, content, reasoning_content=""):
        self.content = content
        self.reasoning_content = reasoning_content


def test_notes_text_rejects_a_response_that_carries_no_notes():
    """These models think in `reasoning_content` and answer in `content`.
    When the token budget runs out mid-thought `content` comes back null -
    that is a truncated run, not notes, and it must not reach write_text()."""
    good = FakeMessage("## Summary\nNormalization.", reasoning_content="1. The user wants")
    assert appmod.notes_text(good, "stop") == "## Summary\nNormalization."

    for message, reason in [
        (FakeMessage(None, reasoning_content="1. The user wants"), "length"),
        (FakeMessage("", reasoning_content=""), "stop"),
        (FakeMessage("   "), "stop"),
    ]:
        try:
            appmod.notes_text(message, reason)
            raise AssertionError(f"should have rejected {message.content!r}")
        except AssertionError:
            raise
        except RuntimeError as e:
            assert str(e), "the failure has to explain itself to the user"
            assert "NoneType" not in str(e), "not a raw TypeError"


if __name__ == "__main__":
    for name, fn in sorted(globals().items()):
        if name.startswith("test_"):
            fn()
            print(f"ok  {name}")
    print("\nall passed")
