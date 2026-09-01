"""Text extraction from class materials, so the notes can be structured around them."""
import json
from pathlib import Path

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".heic", ".webp", ".tiff"}


def from_pdf(path):
    from pypdf import PdfReader
    pages = []
    for i, page in enumerate(PdfReader(str(path)).pages, 1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append(f"[Page {i}]\n{text}")
    return "\n\n".join(pages)


def from_pptx(path):
    from pptx import Presentation
    slides = []
    for i, slide in enumerate(Presentation(str(path)).slides, 1):
        lines = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if lines:
            slides.append(f"[Slide {i}]\n" + "\n".join(lines))
    return "\n\n".join(slides)


def from_docx(path):
    from docx import Document
    return "\n".join(p.text for p in Document(str(path)).paragraphs if p.text.strip())


def from_image(path):
    """On-device OCR via macOS Vision - no model download, no API call."""
    from ocrmac import ocrmac
    lines = ocrmac.OCR(str(path)).recognize()
    return "\n".join(text for text, _conf, _box in lines)


def from_text(path):
    """Markdown and plain text are already text - errors="replace" so a stray
    byte from someone else's editor is not a failed upload."""
    return path.read_text(encoding="utf-8", errors="replace")


EXTRACTORS = {
    ".pdf": from_pdf, ".pptx": from_pptx, ".docx": from_docx,
    ".md": from_text, ".markdown": from_text, ".txt": from_text,
}


def readable_files(directory):
    """Material files in a directory, sorted, skipping dotfiles and unknown types."""
    d = Path(directory)
    if not d.is_dir():
        return []
    out = []
    for p in sorted(d.glob("*")):
        ext = p.suffix.lower()
        if p.is_file() and not p.name.startswith(".") and (ext in EXTRACTORS or ext in IMAGE_EXTS):
            out.append(p)
    return out


def extract_one(path):
    ext = path.suffix.lower()
    fn = EXTRACTORS.get(ext) or (from_image if ext in IMAGE_EXTS else None)
    if not fn:
        return ""
    try:
        return fn(path).strip()
    except Exception as e:
        return f"(could not read {path.name}: {e})"


def extract_all(*dirs, label=None):
    """Concatenate every readable material across dirs, labelled by filename."""
    chunks = []
    for directory in dirs:
        for path in readable_files(directory):
            text = extract_one(path)
            if text:
                tag = f"{label.get(str(directory), '')}{path.name}" if label else path.name
                chunks.append(f"--- {tag} ---\n{text}")
    return "\n\n".join(chunks)


def signature(dirs):
    """Fingerprint of the material set, so cached extraction can be invalidated."""
    parts = []
    for directory in dirs:
        for p in readable_files(directory):
            st = p.stat()
            parts.append(f"{p.name}:{int(st.st_mtime)}:{st.st_size}")
    return "|".join(parts)


def extract_cached(cache_path, *dirs, label=None):
    """Extract, reusing a previous run when the material set is unchanged.

    Called both from the background pre-extract at recording time and from
    note generation, so a class's slides are usually already parsed by the
    time the lecture ends.
    """
    sig = signature(dirs)
    if not sig:
        return ""
    cache = Path(cache_path)
    if cache.exists():
        try:
            cached = json.loads(cache.read_text())
            if cached.get("sig") == sig:
                return cached["text"]
        except (ValueError, KeyError):
            pass  # corrupt cache is just a cache miss
    text = extract_all(*dirs, label=label)
    cache.write_text(json.dumps({"sig": sig, "text": text}))
    return text
